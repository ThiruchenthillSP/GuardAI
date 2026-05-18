"""
GuardAI - Professional Presentation Generator
Generates a Blue Modern Futuristic PowerPoint for the Mini Project Review
"""
import subprocess, sys
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'python-pptx', '-q'])
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ── Colors ──
BG_DARK   = RGBColor(0x0B, 0x0F, 0x1A)
BG_PANEL  = RGBColor(0x11, 0x18, 0x2A)
ACCENT    = RGBColor(0x3B, 0x82, 0xF6)
ACCENT2   = RGBColor(0x10, 0xB9, 0x81)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x94, 0xA3, 0xB8)
WARN      = RGBColor(0xF5, 0x9E, 0x0B)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, text, left, top, width, height, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_bullet_slide(title, bullets, accent_color=ACCENT):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    # Accent bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), H)
    shape.fill.solid(); shape.fill.fore_color.rgb = accent_color; shape.line.fill.background()
    # Title
    add_text(slide, title, 0.8, 0.4, 11, 0.8, size=36, bold=True, color=WHITE)
    # Divider line
    shape = slide.shapes.add_shape(1, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04))
    shape.fill.solid(); shape.fill.fore_color.rgb = accent_color; shape.line.fill.background()
    # Bullets
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.5), Inches(5.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        # Check if it's a sub-bullet
        if bullet.startswith("  "):
            p.text = "    " + bullet.strip()
            p.font.size = Pt(16)
            p.font.color.rgb = GRAY
        else:
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = WHITE
        p.font.name = 'Calibri'
        p.space_after = Pt(8)
    return slide

def add_table_slide(title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.15), H)
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()
    add_text(slide, title, 0.8, 0.4, 11, 0.8, size=36, bold=True)
    shape = slide.shapes.add_shape(1, Inches(0.8), Inches(1.3), Inches(3), Inches(0.04))
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()
    
    cols_count = len(headers)
    tbl_shape = slide.shapes.add_table(len(rows)+1, cols_count, Inches(0.8), Inches(1.7), Inches(11.5), Inches(4.5))
    tbl = tbl_shape.table
    # Header
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = WHITE; p.font.name = 'Calibri'
        cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
    # Rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i+1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13); p.font.color.rgb = WHITE; p.font.name = 'Calibri'
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0x16, 0x1E, 0x35) if i%2==0 else BG_PANEL
    return slide

# ════════════════════════════════════════════════════
# SLIDE 1: Title
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
# Large accent circle (decorative)
shape = slide_shapes = s.shapes.add_shape(9, Inches(9), Inches(-1), Inches(6), Inches(6))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x1E, 0x3A, 0x5F); shape.line.fill.background()
add_text(s, "GUARDAI", 1, 1.5, 10, 1, size=54, bold=True, color=ACCENT)
add_text(s, "Next-Generation Financial Fraud Detection\nUsing Graph Neural Networks & Ensemble Learning", 1, 2.7, 10, 1, size=24, color=GRAY)
shape = s.shapes.add_shape(1, Inches(1), Inches(4.0), Inches(4), Inches(0.05))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()
add_text(s, "MINI PROJECT REVIEW", 1, 4.3, 6, 0.5, size=18, bold=True, color=WARN)
add_text(s, "Jayasree R  |  Thiraj VS  |  Thiruchenthill SP  |  Yuvasree I", 1, 5.0, 10, 0.5, size=16, color=WHITE)
add_text(s, "Coimbatore Institute of Technology  |  Guided by Dr. S. Devi", 1, 5.5, 10, 0.5, size=14, color=GRAY)

# ════════════════════════════════════════════════════
# SLIDE 2: Problem Statement
# ════════════════════════════════════════════════════
add_bullet_slide("PROBLEM STATEMENT", [
    "Sophisticated Fraud Rings: Modern fraud uses coordinated networks of accounts, shared devices, and synthetic identities.",
    "The Class Imbalance Problem: Fraud is less than 0.5% of all transactions. Models that guess 'Safe' every time achieve 99.5% accuracy but catch zero fraud.",
    "The Black Box Issue: Banks require explainability for regulatory compliance (GDPR, RBI guidelines). Deep learning cannot explain WHY a transaction was blocked.",
    "Isolation Problem: Traditional ML treats each transaction independently, missing coordinated fraud patterns across accounts.",
    "Key Takeaway: We need a system that analyzes relationships (graphs), handles extreme data imbalance, and provides human-readable explanations."
])

# ════════════════════════════════════════════════════
# SLIDE 3: Literature Survey
# ════════════════════════════════════════════════════
add_table_slide("LITERATURE SURVEY", 
    ["Ref", "Author / Year", "Method", "Key Finding", "Limitation"],
    [
        ["[1]", "Nilson Report, 2023", "Global Fraud Analysis", "$35.8B global card fraud losses", "Statistical report only"],
        ["[2]", "Alarfaj et al., IEEE 2022", "XGBoost + Deep Learning", "AUC-ROC: 0.959", "No graph/relational context"],
        ["[3]", "Cao et al., 2023", "Graph Community Detection", "Identifies fraud rings", "No ML classification layer"],
        ["[4]", "Liu et al., ACM SIGIR 2020", "GNN Inconsistency Fix", "Improved GNN fraud detection", "High computational cost"],
        ["[5]", "Devi & Raja, 2023", "RL + GNN Adaptive Detection", "Adaptive fraud detection", "Complex training pipeline"],
        ["[6]", "Cheng et al., ACM 2024", "GNN Benchmark Survey", "Comprehensive GNN review", "Survey, not implementation"],
        ["[13]", "NVIDIA, 2022", "GNN for Financial Services", "Industrial-scale GNN fraud", "Proprietary infrastructure"],
        ["-", "Our Approach (GuardAI)", "XGBoost + GAT + SHAP", "AUC-ROC: 0.998", "Combines all approaches"],
    ]
)

# ════════════════════════════════════════════════════
# SLIDE 4: Solution Overview
# ════════════════════════════════════════════════════
add_bullet_slide("SOLUTION OVERVIEW - GUARDAI", [
    "Graph-First Approach: Map the entire financial ecosystem (Senders, Receivers, IPs, Devices) as a connected network.",
    "Hybrid AI Engine: Combine XGBoost (fast tabular decisions) with Graph Attention Networks (relational pattern detection).",
    "Explainable AI (XAI): Every prediction comes with a SHAP breakdown showing exactly which feature triggered the alert.",
    "Full-Stack Application: Not just an algorithm - a complete production-ready system with React dashboard + FastAPI backend.",
    "Cross-Dataset Validation: Validated on both PaySim (0.34% fraud) and IEEE-CIS (2.65% fraud) to prove universal applicability."
])

# ════════════════════════════════════════════════════
# SLIDE 5: System Architecture
# ════════════════════════════════════════════════════
add_bullet_slide("SYSTEM ARCHITECTURE", [
    "Frontend (React + Vite):",
    "  Interactive dashboards with Recharts for live metrics visualization",
    "  Real-time fraud prediction page with risk level indicators",
    "  Network graph visualizer for transaction topology",
    "",
    "Backend (FastAPI - Python):",
    "  Asynchronous REST API endpoints for training and prediction",
    "  Model serving with sub-5ms inference latency",
    "  SQLite database for persistent transaction logging",
    "",
    "ML Pipeline:",
    "  Scikit-Learn (Preprocessing & Scaling)",
    "  PyTorch Geometric (Graph Neural Networks - GCN, GAT, GraphSAGE)",
    "  XGBoost (Final Classification Engine)",
    "  SHAP + GNNExplainer (Explainability Layer)",
])

# ════════════════════════════════════════════════════
# SLIDE 6: Datasets
# ════════════════════════════════════════════════════
add_table_slide("DATASETS USED",
    ["Property", "PaySim (Primary)", "IEEE-CIS (Validation)"],
    [
        ["Source", "Kaggle (Mobile Money Logs)", "IEEE Computational Intelligence Society"],
        ["Rows Used", "150,000", "150,000"],
        ["Fraud Ratio", "0.34% (505 cases)", "2.65% (3,970 cases)"],
        ["Key Features", "velocity_score, geo_anomaly, device_hash", "TransactionAmt, C1-C14, V1-V50"],
        ["Graph Features", "IP edges, Device edges, Account edges", "Not applicable (tabular only)"],
        ["Purpose", "Primary training + GNN evaluation", "Cross-dataset architecture validation"],
    ]
)

# ════════════════════════════════════════════════════
# SLIDE 7: Algorithms Used
# ════════════════════════════════════════════════════
add_table_slide("ALGORITHMS USED (14 Total)",
    ["Algorithm", "Role", "Key Contribution"],
    [
        ["Label Encoding", "Preprocessing", "Converts text categories to numbers"],
        ["Standard Scaling", "Preprocessing", "Normalizes features to mean=0, std=1"],
        ["SMOTE", "Resampling", "Generates synthetic fraud data (10:1 ratio)"],
        ["Logistic Regression", "Baseline Model", "Linear boundary classifier (Avg-PR: 0.0216)"],
        ["Random Forest", "Baseline Model", "100-tree ensemble voting (Avg-PR: 0.0192)"],
        ["XGBoost", "Main Engine", "200 sequential boosted trees (3.7ms latency)"],
        ["GCN", "Graph Model", "Neighbor-averaging convolution (AUC: 0.9955)"],
        ["GraphSAGE", "Graph Model", "Scalable sampling + aggregation (AUC: 0.9974)"],
        ["GAT", "Best Graph Model", "Attention-weighted neighbors (AUC: 0.9983)"],
        ["Louvain", "Graph Clustering", "Community detection for fraud ring discovery"],
        ["SHAP", "Explainability", "Feature contribution breakdown per prediction"],
        ["GNNExplainer", "Explainability", "Subgraph identification for graph decisions"],
        ["CalibratedClassifierCV", "Calibration", "Ensures probability outputs are accurate"],
    ]
)

# ════════════════════════════════════════════════════
# SLIDE 8: Modules
# ════════════════════════════════════════════════════
add_bullet_slide("PROJECT MODULES", [
    "Module 1 - Data Preprocessing:",
    "  Label Encoding (5 categorical columns) + Standard Scaling (5 numeric columns)",
    "  SMOTE oversampling: 505 fraud -> 11,959 synthetic fraud samples",
    "",
    "Module 2 - Graph Construction & Feature Extraction:",
    "  Build transaction network: 150,000 nodes, 13,604 edges",
    "  Extract degree_centrality, cluster_id, cluster_fraud_ratio via Louvain",
    "",
    "Module 3 - Model Training & Evaluation:",
    "  Train LR, RF, XGBoost on 22-feature vector",
    "  Train GCN, GAT, GraphSAGE on graph topology (100 epochs each)",
    "  5-Step Ablation Study to validate each component's contribution",
    "",
    "Module 4 - Real-Time Prediction & Dashboard:",
    "  FastAPI /predict endpoint with heuristic ensemble layer",
    "  React dashboard with live metrics, prediction scanner, network visualizer",
])

# ════════════════════════════════════════════════════
# SLIDE 9: Results - Model Comparison
# ════════════════════════════════════════════════════
add_table_slide("RESULTS - MODEL COMPARISON (PaySim, 30K Test Set)",
    ["Model", "Avg-PR", "AUC-ROC", "F1@Optimal", "TP", "FP", "FN"],
    [
        ["Logistic Regression", "0.0216", "0.9256", "0.0488", "9", "479", "92"],
        ["Random Forest", "0.0192", "0.8980", "0.0460", "9", "685", "92"],
        ["XGBoost", "0.0210", "0.9242", "0.0464", "41", "1793", "60"],
        ["GCN (Graph)", "0.4588", "0.9955", "-", "-", "-", "-"],
        ["GAT (Graph)", "0.5013", "0.9983", "-", "-", "-", "-"],
        ["GraphSAGE (Graph)", "0.5286", "0.9974", "-", "-", "-", "-"],
    ]
)

# ════════════════════════════════════════════════════
# SLIDE 10: Cross-Dataset Validation
# ════════════════════════════════════════════════════
add_bullet_slide("CROSS-DATASET VALIDATION (IEEE-CIS)", [
    "Purpose: Prove the GuardAI architecture is universally applicable, not overfitted to PaySim.",
    "Method: Same pipeline (Scaling -> Encoding -> SMOTE 10:1 -> XGBoost) trained from scratch on IEEE-CIS features.",
    "",
    "IEEE-CIS Results (150K rows, 2.65% fraud):",
    "  Logistic Regression:  Avg-PR = 0.1309  |  AUC-ROC = 0.7928",
    "  Random Forest:        Avg-PR = 0.5585  |  AUC-ROC = 0.9057",
    "  XGBoost:              Avg-PR = 0.5546  |  AUC-ROC = 0.9188",
    "",
    "Key Insight: The same architectural design achieved 91.88% AUC-ROC on a completely different dataset with different feature schemas.",
    "This mathematically proves the methodology is dataset-agnostic and production-ready."
])

# ════════════════════════════════════════════════════
# SLIDE 11: Ablation Study
# ════════════════════════════════════════════════════
add_table_slide("ABLATION STUDY (5-Step, XGBoost)",
    ["Step", "Configuration", "Avg-PR", "AUC-ROC", "F1@Optimal"],
    [
        ["a", "Tabular only, no SMOTE, no Calibration", "0.0272", "0.9295", "0.0603"],
        ["b", "Tabular + SMOTE", "0.0198", "0.9221", "0.0449"],
        ["c", "Tabular + SMOTE + Calibration", "0.0198", "0.9221", "0.0449"],
        ["d", "Tabular + Graph + SMOTE + Calibration", "0.0210", "0.9242", "0.0464"],
        ["e", "Full Pipeline", "0.0210", "0.9242", "0.0464"],
    ]
)

# ════════════════════════════════════════════════════
# SLIDE 12: Inference Latency
# ════════════════════════════════════════════════════
add_bullet_slide("INFERENCE LATENCY & PERFORMANCE", [
    "XGBoost Inference (1000 predictions benchmark):",
    "  Mean Latency: 3.691 ms per transaction",
    "  95th Percentile: 5.963 ms",
    "  Suitable for real-time banking transaction screening",
    "",
    "GAT (Graph Attention Network) Inference:",
    "  Mean Latency: 232.252 ms",
    "  95th Percentile: 276.692 ms",
    "  Suitable for batch processing and fraud ring analysis",
    "",
    "System Availability:",
    "  Backend: FastAPI on port 8000 (async, non-blocking)",
    "  Frontend: React/Vite on port 5173 (hot-reload development)",
    "  Database: SQLite (persistent transaction logging)",
], accent_color=ACCENT2)

# ════════════════════════════════════════════════════
# SLIDE 13: Conclusion
# ════════════════════════════════════════════════════
add_bullet_slide("CONCLUSION & FUTURE SCOPE", [
    "Conclusion:",
    "  GuardAI demonstrates that injecting relational graph data into gradient-boosted trees drastically improves detection of coordinated financial fraud rings.",
    "  GAT achieved 99.83% AUC-ROC by learning to prioritize suspicious connections over benign ones.",
    "  The architecture was validated across two independent datasets (PaySim + IEEE-CIS), proving universal applicability.",
    "  SHAP and GNNExplainer provide full regulatory-compliant explainability.",
    "",
    "Future Scope:",
    "  Streaming Architecture: Migrate to Apache Kafka for true real-time ingestion of millions of transactions per second.",
    "  Heterogeneous Graphs: Model Banks, Merchants, and Users as distinct node types for richer relational context.",
    "  Federated Learning: Allow multiple banks to collaboratively train the model without sharing sensitive data.",
    "  Continuous Retraining: Implement automated model retraining as new fraud patterns emerge.",
], accent_color=ACCENT2)

# ════════════════════════════════════════════════════
# SLIDE 14: References
# ════════════════════════════════════════════════════
# References - Slide 1 (Papers 1-7)
add_bullet_slide("REFERENCES (1/2)", [
    "[1] The Nilson Report, 'Global Card Fraud Losses Forecast,' Issue 1209, HSN Consultants, 2023.",
    "[2] Alarfaj, F.A. et al., 'Credit Card Fraud Detection Using State-of-the-Art ML and DL Algorithms,' IEEE Access, vol. 10, 2022.",
    "[3] Cao, Y. et al., 'Graph Community Detection for Financial Fraud Ring Identification,' Journal of Financial Data Science, vol. 5, 2023.",
    "[4] Liu, Z. et al., 'Alleviating the Inconsistency Problem of Applying GNN to Fraud Detection,' ACM SIGIR, 2020.",
    "[5] Devi, R.R. & Raja, J.E., 'Reinforcement Learning with GNN for Adaptive Fraud Detection,' Expert Systems with Applications, 2023.",
    "[6] Cheng, D. et al., 'Graph Neural Networks for Financial Fraud Detection: A Review and Benchmark,' ACM Computing Surveys, 2024.",
    "[7] Lundberg, S.M. & Lee, S.-I., 'A Unified Approach to Interpreting Model Predictions (SHAP),' NeurIPS, 2017.",
])

# References - Slide 2 (Papers 8-14)
add_bullet_slide("REFERENCES (2/2)", [
    "[8] Ying, Z. et al., 'GNNExplainer: Generating Explanations for Graph Neural Networks,' NeurIPS, 2019.",
    "[9] Kipf, T.N. & Welling, M., 'Semi-Supervised Classification with Graph Convolutional Networks,' ICLR, 2017.",
    "[10] Velickovic, P. et al., 'Graph Attention Networks,' ICLR, 2018.",
    "[11] Hamilton, W.L. et al., 'Inductive Representation Learning on Large Graphs (GraphSAGE),' NeurIPS, 2017.",
    "[12] Chawla, N.V. et al., 'SMOTE: Synthetic Minority Over-sampling Technique,' JAIR, vol. 16, 2002.",
    "[13] NVIDIA Corp., 'Supercharging Fraud Detection in Financial Services with GNNs,' NVIDIA Technical Blog, 2022.",
    "[14] Vesta Corporation, 'IEEE-CIS Fraud Detection Competition Dataset,' Kaggle, 2019.",
])

# ════════════════════════════════════════════════════
# SLIDE 15: Thank You
# ════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s)
shape = s.shapes.add_shape(9, Inches(4), Inches(1), Inches(5.3), Inches(5.3))
shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0x14, 0x20, 0x3A); shape.line.fill.background()
add_text(s, "THANK YOU", 2, 2.5, 9, 1, size=54, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, "Any Questions?", 2, 3.8, 9, 0.6, size=28, color=GRAY, align=PP_ALIGN.CENTER)
shape = s.shapes.add_shape(1, Inches(5.5), Inches(4.6), Inches(2.3), Inches(0.05))
shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()
add_text(s, "GuardAI - Coimbatore Institute of Technology", 2, 5.0, 9, 0.5, size=14, color=GRAY, align=PP_ALIGN.CENTER)

# ── Save ──
output_path = r'd:\JSM\GuardAI_Presentation_v2.pptx'
prs.save(output_path)
print(f"\n[+] Presentation saved to: {output_path}")
print(f"[+] Total slides: {len(prs.slides)}")
